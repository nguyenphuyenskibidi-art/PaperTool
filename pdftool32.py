import fitz  # PyMuPDF
from pptx import Presentation
import os
import pikepdf
import sys

# --- TỪ ĐIỂN NGÔN NGỮ ---
LANGS = {
    'vn': {
        'select_lang': "CHỌN NGÔN NGỮ: 1. Tiếng Việt | 2. English: ",
        'menu_title': "💎 PDF-PPTX",
        'opt1': "1. Chuyển PDF -> PPTX (Fix Font & High-DPI)",
        'opt2': "2. Phân tích chi tiết (Dòng, Chữ, Ảnh, Metadata)",
        'opt3': "3. Quản lý (Ghép, Tách, Cắt, Đặt mật khẩu)",
        'opt4': "4. Thoát",
        'ask_choice': "\n[?] Nhập Lệnh  (1-4): ",
        'ask_file': "> Kéo thả file PDF vào đây: ",
        'err_file': "[!] Lỗi: File không tồn tại!",
        'back_menu': "XỬ LÝ XONG! NHẤN ENTER ĐỂ QUAY LẠI MENU...",
        'proc_analyze': "[🔍] ĐANG PHÂN TÍCH SÂU...",
        'proc_convert': "[🚀] ĐANG RENDER  (300 DPI)...",
        'done': "[💎] THÀNH CÔNG!",
        'm_title': "\n--- QUẢN LÝ CẤU TRÚC PDF ---",
        'm_1': "1. Ghép PDF", 'm_2': "2. Tách/Cắt trang", 'm_3': "3. Đặt mật khẩu", 'm_4': "4. Quay lại"
    },
    'en': {
        'select_lang': "SELECT LANGUAGE: 1. Tiếng Việt | 2. English: ",
        'menu_title': "💎 PDF-PPTX",
        'opt1': "1. Convert PDF -> PPTX (Fix Font & High-DPI)",
        'opt2': "2. Deep Analysis (Lines, Words, Images, Metadata)",
        'opt3': "3. Management (Merge, Split, Crop, Password)",
        'opt4': "4. Exit",
        'ask_choice': "\n[?] Your choice (1-4): ",
        'ask_file': "> Drag and drop PDF file here: ",
        'err_file': "[!] Error: File not found!",
        'back_menu': "PROCESS DONE! PRESS ENTER TO RETURN TO MENU...",
        'proc_analyze': "[🔍] DEEP ANALYZING...",
        'proc_convert': "[🚀] RENDERING (300 DPI)...",
        'done': "[💎] SUCCESS!",
        'm_title': "\n--- PDF STRUCTURE MANAGEMENT ---",
        'm_1': "1. Merge PDF", 'm_2': "2. Split/Crop Pages", 'm_3': "3. Set Password", 'm_4': "4. Back"
    }
}

# --- CÁC HÀM XỬ LÝ CHÍNH ---
def analyze_deep(path, L):
    print(f"\n{L['proc_analyze']}")
    try:
        doc = fitz.open(path)
        print(f"[*] Metadata: {doc.metadata}")
        for i, page in enumerate(doc):
            words = len(page.get_text("words"))
            imgs = len(page.get_images())
            lines = len(page.get_text().splitlines())
            print(f"Slide {i+1}: {words} chữ | {lines} dòng | {imgs} ảnh")
        doc.close()
    except Exception as e: print(f"Error: {e}")

def convert_to_pptx(path, L):
    try:
        output = path.rsplit(".", 1)[0] + "_Tool_.pptx"
        prs = Presentation()
        doc = fitz.open(path)
        print(f"\n{L['proc_convert']}")
        for i in range(len(doc)):
            page = doc.load_page(i)
            # Render Matrix(3,3) để fix lỗi font và làm mượt hình ảnh
            pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))
            img_temp = f"tmp_{i}.png"
            pix.save(img_temp)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_temp, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(img_temp)
            print(f"  > Page {i+1}/{len(doc)} OK")
        prs.save(output)
        print(f"\n{L['done']} -> {output}")
        doc.close()
    except Exception as e: print(f"Error: {e}")

def manage_pdf(path, L):
    while True:
        print(L['m_title'])
        print(f"{L['m_1']}\n{L['m_2']}\n{L['m_3']}\n{L['m_4']}")
        m_choice = input("> ")
        if m_choice == '4': break
        try:
            if m_choice == '1':
                f2 = input("File 2: ").strip().replace('"', '')
                with pikepdf.open(path) as p1, pikepdf.open(f2) as p2:
                    p1.pages.extend(p2.pages)
                    p1.save(path.replace(".pdf", "_Merged.pdf"))
                print(L['done'])
            elif m_choice == '2':
                s, e = int(input("Start: ")), int(input("End: "))
                with pikepdf.open(path) as p:
                    new = pikepdf.Pdf.new()
                    new.pages.extend(p.pages[s-1:e])
                    new.save(path.replace(".pdf", "_Split.pdf"))
                print(L['done'])
            elif m_choice == '3':
                pw = input("Password: ")
                with pikepdf.open(path) as p:
                    p.save(path.replace(".pdf", "_Locked.pdf"), encryption=pikepdf.Encryption(user=pw, owner=pw))
                print(L['done'])
        except Exception as e: print(f"Error: {e}")

# --- MAIN LOOP ---
if __name__ == "__main__":
    os.system('chcp 65001 > nul')
    l_idx = input(LANGS['vn']['select_lang'])
    L = LANGS['vn'] if l_idx == '1' else LANGS['en']

    while True:
        os.system('cls' if os.name == 'nt' else 'clear')
        print("="*50 + f"\n      {L['menu_title']}\n" + "="*50)
        print(f"{L['opt1']}\n{L['opt2']}\n{L['opt3']}\n{L['opt4']}")
        
        choice = input(L['ask_choice'])
        if choice == '4': break
        
        file_p = input(L['ask_file']).strip().replace('"', '').replace("'", "")
        if not os.path.exists(file_p):
            print(L['err_file']); input("Enter..."); continue

        if choice == '1': convert_to_pptx(file_p, L)
        elif choice == '2': analyze_deep(file_p, L)
        elif choice == '3': manage_pdf(file_p, L)
        
        input("\n" + L['back_menu'])